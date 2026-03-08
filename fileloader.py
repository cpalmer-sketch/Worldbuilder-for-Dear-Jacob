import os
import csv
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
import pandas as pd

def extract_txt(path):
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()

def extract_csv(path):
    with open(path, 'r', encoding='utf-8') as f:
        reader = csv.reader(f)
        return '\n'.join([', '.join(row) for row in reader])

def extract_docx(path):
    doc = Document(path)
    return '\n'.join([para.text for para in doc.paragraphs])

def extract_xlsx(path):
    df = pd.read_excel(path)
    return df.to_string(index=False)

def extract_pdf(path):
    reader = PdfReader(path)
    text = ''
    for page in reader.pages:
        if page.extract_text():
            text += page.extract_text() + '\n'
    return text

def extract_pptx(path):
    prs = Presentation(path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def extract_file(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == '.txt':
        return extract_txt(path)
    elif ext == '.csv':
        return extract_csv(path)
    elif ext == '.docx':
        return extract_docx(path)
    elif ext == '.xlsx':
        return extract_xlsx(path)
    elif ext == '.pdf':
        return extract_pdf(path)
    elif ext == '.pptx':
        return extract_pptx(path)
    else:
        return f'Unsupported file format: {ext}'
